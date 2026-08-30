import sys
import os
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt
# pyrefly: ignore [missing-import]
from pptx.enum.text import PP_ALIGN
# pyrefly: ignore [missing-import]
from pptx.dml.color import RGBColor
# pyrefly: ignore [missing-import]
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(filename="ControlPlane_Business_Proposal.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    DARK_BG = RGBColor(15, 23, 42)        # #0F172A (Deep Slate / Dark Mode)
    LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC (Clean Light Slate)
    ACCENT_PURPLE = RGBColor(124, 58, 237)# #7C3AED (Accenture/Modern Purple)
    ACCENT_BLUE = RGBColor(14, 165, 233)  # #0EA5E9 (Cyan Blue)
    TEXT_DARK = RGBColor(30, 41, 59)      # #1E293B
    TEXT_LIGHT = RGBColor(241, 245, 249)  # #F1F5F9
    MUTED_TEXT = RGBColor(100, 116, 139)  # #64748B
    CARD_BG = RGBColor(255, 255, 255)     # White card
    CARD_BORDER = RGBColor(226, 232, 240) # Slate 200

    def add_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="ACCENTURE INNOVATION CHALLENGE 2026"):
        # Category / Kicker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_PURPLE
        p_cat.font.name = "Calibri"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.font.name = "Calibri"

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_background(s1, DARK_BG)

    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "ControlPlane.ai"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE
    p1.font.name = "Calibri"

    p2 = tf1.add_paragraph()
    p2.text = "Enterprise AI Governance, Responsible Agent Security & Safety Gateway"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = "Calibri"
    p2.space_before = Pt(10)

    sub_box = s1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.2))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Submission Track: Accenture Innovation Challenge 2026 — Prototype Development\nRepository: github.com/MohithChandra07/Control-Plane-AIC"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = MUTED_TEXT
    p_sub.font.name = "Calibri"

    # ==========================================
    # SLIDE 2: Executive Summary & Problem
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_background(s2, LIGHT_BG)
    add_header(s2, "Executive Summary: The Enterprise AI Trust Gap")

    # Card 1: Problem
    add_card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Problem & Vulnerabilities"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    bullets_p = [
        "Unmonitored LLM Output: Hallucinated claims cause direct financial loss (e.g. false refund promises).",
        "Rogue Tool Execution: Tainted conversational claims trigger downstream transactional API calls (money transfers, DB overwrites).",
        "Data Leakage & PII: Exposure of sensitive user info violating GDPR/HIPAA.",
        "Regulatory Pressure: EU AI Act & NIST AI RMF mandate verifiable audit trails and human oversight."
    ]
    for b in bullets_p:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(10)

    # Card 2: Solution
    add_card(s2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb2 = s2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "The ControlPlane Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    bullets_s = [
        "Inline Governance Gateway: Sits transparently between applications and LLMs (OpenAI compatible).",
        "Adaptive Scrutiny: Multi-tier inspection balancing cost, latency, and rigorous verification.",
        "Multi-Turn Taint Tracking: Blocks tool calls derived from unverified earlier conversation claims.",
        "Surgical Remediation: Redacts, hedges, or cites per-claim rather than blocking entire responses.",
        "Immutable Audit Ledger: Cryptographically hash-chained transaction logs for complete compliance."
    ]
    for b in bullets_s:
        p_b = tf2.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: System Architecture Flow
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_background(s3, LIGHT_BG)
    add_header(s3, "System Architecture: End-to-End Governance Gateway")

    cards_arch = [
        ("1. Request Interception", "OpenAI-compatible drop-in gateway endpoint (/v1/chat/completions) with tenant header routing.", ACCENT_PURPLE),
        ("2. Adaptive Scrutiny", "Tier 0 cheap heuristic pre-filter → Tier 1 claim verification & regex PII scanner.", ACCENT_BLUE),
        ("3. Policy & Remediation", "Per-claim surgical remediation (Hedge, Redact, Escalate) & live risk appetite slider.", ACCENT_PURPLE),
        ("4. Taint & Tool Gating", "Intercepts high-risk tool calls (e.g. issue_refund) if inputs rely on unverified claims.", ACCENT_BLUE)
    ]

    for i, (title, desc, color) in enumerate(cards_arch):
        x = Inches(0.8 + i * 3.0)
        add_card(s3, x, Inches(1.8), Inches(2.7), Inches(4.8))
        tb = s3.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.4), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_before = Pt(12)

    # ==========================================
    # SLIDE 4: Empirical Results & Verification
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_background(s4, LIGHT_BG)
    add_header(s4, "Empirical Validation & Benchmark Results")

    stats = [
        ("145 / 145", "Automated Tests Passing", "100% test suite pass rate covering all integration scenes.", ACCENT_PURPLE),
        ("10,000", "Synthetic Interactions", "Replayed in ~100s, verifying high-throughput audit ledger.", ACCENT_BLUE),
        ("< 0.05", "ECE Calibration Error", "Expected Calibration Error measured on synthetic benchmark.", ACCENT_PURPLE),
        ("3 Policies", "Per-Tenant Multi-Tenancy", "Different outcomes enforced for Support, Copilot, Regulated.", ACCENT_BLUE)
    ]

    for i, (num, label, detail, color) in enumerate(stats):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.6)

        add_card(s4, x, y, Inches(5.6), Inches(2.3))
        tb = s4.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.0), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

        p_l = tf.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(14)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_DARK

        p_det = tf.add_paragraph()
        p_det.text = detail
        p_det.font.size = Pt(11)
        p_det.font.color.rgb = MUTED_TEXT
        p_det.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: Business Model & Market Growth
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_background(s5, LIGHT_BG)
    add_header(s5, "Market Opportunity & Business Model")

    # Left Card: Market & Strategy
    add_card(s5, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "$35B Enterprise AI TRiSM TAM"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    b_m = [
        "SaaS Gateway Pricing: Token-volume base + monthly managed tenant subscription.",
        "Enterprise License: On-Prem/VPC flat license per gateway node.",
        "Accenture Synergy: Integrated asset for Accenture AI Advisory & Transformation projects.",
        "Target Verticals: Banking, Healthcare, Customer Service, Public Sector."
    ]
    for b in b_m:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(10)

    # Right Card: Financial Trajectory
    add_card(s5, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb2 = s5.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Financial Forecast (3-Year)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    fin_lines = [
        "Year 1: 15 Enterprise Clients | $1.2M ARR (78% Gross Margin)",
        "Year 2: 60 Enterprise Clients | $5.8M ARR (82% Gross Margin)",
        "Year 3: 220 Enterprise Clients | $24.5M ARR (86% Gross Margin)",
        "Key Cost Drivers: Low infrastructure overhead due to lightweight Tier 0 pre-filtering."
    ]
    for b in fin_lines:
        p_b = tf2.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(12)

    # Save Presentation
    prs.save(filename)
    print(f"Successfully generated {filename}")

if __name__ == "__main__":
    create_presentation()
